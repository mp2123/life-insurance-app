import os
import re
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt

def convert_to_word(md_path, docx_path):
    doc = Document()
    doc.add_heading('Life Insurance AZ 2026: Master Study Guide', 0)
    
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
        
    blocks = content.split('\n---')
    for block in blocks:
        if '❓ Question' in block:
            q_match = re.search(r'## ❓ Question (\d+)', block)
            if q_match:
                doc.add_heading(f'Question {q_match.group(1)}', level=1)
            
            q_text_match = re.search(r'\*\*(.*?)\*\*', block, re.DOTALL)
            if q_text_match:
                p = doc.add_paragraph()
                p.add_run(q_text_match.group(1).strip()).bold = True
            
            if '### 📝 Options' in block:
                doc.add_heading('Options', level=2)
                options_section = block.split('### 📝 Options')[1].split('### 💡 Explanation')[0]
                for line in options_section.split('\n'):
                    line = line.strip()
                    if line.startswith('-'):
                        p = doc.add_paragraph(style='List Bullet')
                        if '[x]' in line:
                            run = p.add_run("✅ " + line.replace('- [x]', '').replace('**', '').replace('✅', '').strip())
                            run.bold = True
                        else:
                            p.add_run(line.replace('- [ ]', '').strip())
            
            if '### 💡 Explanation' in block:
                doc.add_heading('Explanation', level=2)
                expl = block.split('### 💡 Explanation')[1].strip().strip('>')
                p = doc.add_paragraph(expl)
                p.italic = True

    doc.save(docx_path)

def convert_to_powerpoint(md_path, pptx_path):
    prs = Presentation()
    
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
        
    blocks = content.split('\n---')
    for block in blocks:
        if '❓ Question' not in block: continue
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        q_match = re.search(r'## ❓ Question (\d+)', block)
        title = f"Question {q_match.group(1)}" if q_match else "Study Question"
        slide.shapes.title.text = title
        
        q_text_match = re.search(r'\*\*(.*?)\*\*', block, re.DOTALL)
        if q_text_match:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = q_text_match.group(1).strip()
            
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"{title}: Answer & Explanation"
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        if '### 📝 Options' in block:
            options_section = block.split('### 📝 Options')[1].split('### 💡 Explanation')[0]
            for line in options_section.split('\n'):
                line = line.strip()
                if line.startswith('-'):
                    p = tf.add_paragraph()
                    if '[x]' in line:
                        p.text = "✅ " + line.replace('- [x]', '').replace('**', '').replace('✅', '').strip()
                        p.font.bold = True
                    else:
                        p.text = "☐ " + line.replace('- [ ]', '').strip()
        
        if '### 💡 Explanation' in block:
            p = tf.add_paragraph()
            p.text = "\nExplanation:"
            p.font.bold = True
            expl = block.split('### 💡 Explanation')[1].strip().strip('>')
            p = tf.add_paragraph()
            p.text = (expl[:250] + '...') if len(expl) > 250 else expl
            p.font.size = PptPt(14)
            p.font.italic = True

    prs.save(pptx_path)

if __name__ == "__main__":
    base_dir = "/Users/michael_s_panico/Desktop/Life insurance photos/ALL OCR OCR extracted images for Life Insurance AZ 2026"
    input_file = os.path.join(base_dir, "Life_Insurance_AZ_2026_Study_Guide_Polished.md")
    word_file = os.path.join(base_dir, "Life_Insurance_Study_Guide.docx")
    ppt_file = os.path.join(base_dir, "Life_Insurance_Study_Guide.pptx")
    convert_to_word(input_file, word_file)
    convert_to_powerpoint(input_file, ppt_file)
