import re
import os
from docx import Document
from pptx import Presentation
from pptx.util import Pt as PptPt

LAWS_OF_INSURANCE = {
    "four times": "waiting period of one year",
    "mortgage": "Decreasing term",
    "share in a company's profits": "Profit-Sharing Plan",
    "non-smoker": "Mortality and lifestyle risks",
    "employer owns": "Group life insurance",
    "verifying that a license applicant": "Director of the Department of Insurance",
    "Sarah applies": "June 1",
    "MIB": "Medical history information shared among insurers",
    "Adverse selection": "Insuring high-risk individuals more frequently",
    "Human Life Value": "earnings",
    "7-Pay test": "Modified Endowment Contract (MEC)",
    "partial withdrawals": "Universal life",
    "Convertible Term": "without proof of insurability",
    "last surviving": "Survivorship life",
    "Juvenile": "Payor Provision",
    "Face amount minus cash value": "Net amount at risk",
}

def ultra_clean(text):
    text = re.sub(r'[©®™]', '', text)
    text = re.sub(r' \. \. \. .*?\.com', '', text)
    text = re.sub(r'\d+ of \d+ Questions Remaining', '', text)
    text = re.sub(r'Proceed —?_?>', '', text)
    text = re.sub(r'Submit Exam.*', '', text)
    text = re.sub(r'^\d+:\d+ (?:am|pm|al|ef|wo|oll).*', '', text, flags=re.MULTILINE)
    text = re.sub(r'(\n\s*)+', '\n', text).strip()
    return text

def build_ultimate_guide():
    base_path = "/Users/michael_s_panico/Desktop/Life insurance photos"
    master_data = []
    seen_ids = set()

    for folder_num in range(1, 8):
        folder_path = os.path.join(base_path, str(folder_num))
        if not os.path.exists(folder_path): continue
        
        files = sorted([f for f in os.listdir(folder_path) if f.endswith(".txt")])
        for filename in files:
            try:
                with open(os.path.join(folder_path, filename), 'r', encoding='utf-8') as f:
                    content = f.read()
            except: continue
            
            q_match = re.search(r'Question (\d+).*?\n(.*?)(?=\n[xgiv©*O]|\n\d+ of|\nProceed|\nSubmit|$)', content, re.DOTALL | re.IGNORECASE)
            if not q_match: continue
            
            q_text = ultra_clean(q_match.group(2).strip())
            q_id = re.sub(r'\W+', '', q_text.lower())
            if q_id in seen_ids or len(q_id) < 10: continue
            seen_ids.add(q_id)

            options_raw = re.findall(r'\n([xgiv©*O].*?)(?=\n[xgiv©*O]|\nProceed|\n\d+ of|\nSubmit|$)', content, re.DOTALL)
            options = [re.sub(r'^[xgiv©*O)]+\s*', '', o).strip() for o in options_raw]
            options = [o for o in options if len(o) > 1]

            expl_match = re.search(r'(?:ds|és|is)\s+(.*)', content, re.DOTALL | re.IGNORECASE)
            explanation = ultra_clean(expl_match.group(1).strip()) if expl_match else ""

            correct_ans = None
            for key, val in LAWS_OF_INSURANCE.items():
                if key.lower() in q_text.lower():
                    for opt in options:
                        if val.lower() in opt.lower():
                            correct_ans = opt; break
                    if correct_ans: break
            
            if not correct_ans and explanation:
                for opt in options:
                    if opt.lower() in explanation.lower() and len(opt) > 5:
                        correct_ans = opt; break
            
            if not correct_ans:
                for opt_raw in options_raw:
                    if any(opt_raw.lower().startswith(s) for s in ['g', 'iv', 'v', '©', '*']):
                        correct_ans = re.sub(r'^[xgiv©*O)]+\s*', '', opt_raw).strip()
                        break

            master_data.append({"q": q_text, "opts": options, "ans": correct_ans, "expl": explanation})

    out_dir = os.path.join(base_path, "ALL OCR OCR extracted images for Life Insurance AZ 2026")
    
    with open(os.path.join(out_dir, "ULTIMATE_AZ_2026_STUDY_GUIDE.md"), 'w') as f:
        f.write("# 🎓 ULTIMATE Life Insurance AZ 2026 Study Guide\n\n")
        f.write(f"> **Comprehensive Count:** {len(master_data)} Verified Questions\n\n---\n\n")
        for i, data in enumerate(master_data):
            f.write(f"## ❓ Q{i+1}: {data['q']}\n\n")
            if not data['opts']: f.write("_Options messy in source_\n")
            for o in data['opts']:
                f.write(f"- {'✅ [x] **' + o + '**' if o == data['ans'] else '    [ ] ' + o}\n")
            if data['expl']: f.write(f"\n> 💡 **Explanation:** {data['expl']}\n")
            f.write("\n---\n\n")

    prs = Presentation()
    for i, data in enumerate(master_data):
        s1 = prs.slides.add_slide(prs.slide_layouts[1])
        s1.shapes.title.text = f"Q{i+1}"; s1.placeholders[1].text = data['q']
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = f"A{i+1}"; tf = s2.placeholders[1].text_frame
        for o in data['opts']:
            p = tf.add_paragraph(); p.text = f"{'✅' if o == data['ans'] else '☐'} {o}"
            if o == data['ans']: p.font.bold = True
        if data['expl']:
            p = tf.add_paragraph(); p.text = f"\nNote: {data['expl'][:250]}..."; p.font.size = PptPt(12); p.font.italic = True
    prs.save(os.path.join(out_dir, "ULTIMATE_STUDY_DECK.pptx"))

    doc = Document()
    doc.add_heading('Ultimate Study Guide', 0)
    for i, data in enumerate(master_data):
        doc.add_heading(f'Q{i+1}: {data["q"]}', level=1)
        for o in data['opts']:
            p = doc.add_paragraph(style='List Bullet')
            if o == data['ans']: p.add_run(f"✅ {o}").bold = True
            else: p.add_run(o)
    doc.save(os.path.join(out_dir, "ULTIMATE_STUDY_GUIDE.docx"))

if __name__ == "__main__":
    build_ultimate_guide()
