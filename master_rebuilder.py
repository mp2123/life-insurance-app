import re
import os
from docx import Document
from pptx import Presentation
from pptx.util import Pt as PptPt

# Enhanced knowledge base with exact phrasing found in OCR
KNOWLEDGE_BASE = {
    "Sarah applies": "June 1",
    "Medical Information Bureau": "Medical history information shared among insurers",
    "Adverse selection": "Insuring high-risk individuals more frequently",
    "Human Life Value": "future earnings and contributions",
    "Automatic Premium Loan": "non-payment",
    "larger numbers of similar risks": "Law of large numbers",
    "backdating": "Setting the effective date of the policy earlier than the present",
    "Executive Bonus Plan": "Premiums",
    "Key employee insurance": "Financial stability",
    "failing the insurance licensing exam four times": "waiting period of one year",
    "surrender charge": "early withdrawal",
    "nonforfeiture": "cash value",
    "reinstatement": "unpaid premiums",
    "Incontestability": "two years",
    "Suicide clause": "two years",
    "Grace period": "31 days"
}

def clean_ocr(text):
    text = re.sub(r'[©®™]', '', text)
    text = re.sub(r' \. \. \. .*?\.com', '', text)
    text = re.sub(r'\d+ of \d+ Questions Remaining', '', text)
    text = re.sub(r'Proceed —?_?>', '', text)
    text = re.sub(r'Submit Exam.*', '', text)
    text = re.sub(r'I< >I.*', '', text)
    text = re.sub(r'% 4a.*', '', text)
    text = re.sub(r'< \+ .*', '', text)
    text = re.sub(r'^\d+:\d+ (?:am|pm|al|ef|wo|oll).*', '', text, flags=re.MULTILINE)
    text = re.sub(r'(\n\s*)+', '\n', text).strip()
    return text

def rebuild():
    base_dir = "/Users/michael_s_panico/Desktop/Life insurance photos/ALL OCR OCR extracted images for Life Insurance AZ 2026"
    raw_path = os.path.join(base_dir, "consolidated_raw.txt")
    
    with open(raw_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    chunks = content.split('--- File: ')
    questions = []
    seen = set()

    for chunk in chunks:
        if not chunk.strip(): continue
        
        q_match = re.search(r'Question (\d+).*?\n(.*?)(?=\n[xgiv©*O]|\n\d+ of|\nProceed|\nSubmit|$)', chunk, re.DOTALL | re.IGNORECASE)
        if not q_match: continue
        
        q_text = clean_ocr(q_match.group(2).strip())
        q_id = re.sub(r'\W+', '', q_text.lower())
        if q_id in seen or len(q_id) < 5: continue
        seen.add(q_id)

        # Options Extraction - Improved to handle multi-line options
        options_raw = re.findall(r'\n([xgiv©*O].*?)(?=\n[xgiv©*O]|\nProceed|\n\d+ of|\nSubmit|$)', chunk, re.DOTALL)
        options = []
        for o in options_raw:
            clean_opt = re.sub(r'^[xgiv©*O)]+\s*', '', o).strip()
            # Clean up broken lines within options
            clean_opt = re.sub(r'\s+', ' ', clean_opt)
            if clean_opt and len(clean_opt) > 1:
                options.append(clean_opt)

        expl_match = re.search(r'(?:ds|és|is)\s+(.*)', chunk, re.DOTALL | re.IGNORECASE)
        explanation = clean_ocr(expl_match.group(1).strip()) if expl_match else ""

        # Logic to find correct answer
        correct_ans = None
        for key, ans in KNOWLEDGE_BASE.items():
            if key.lower() in q_text.lower():
                for opt in options:
                    if ans.lower() in opt.lower():
                        correct_ans = opt
                        break
                if correct_ans: break
        
        if not correct_ans and explanation:
            for opt in options:
                # Use a similarity/containment check
                if opt.lower() in explanation.lower() and len(opt) > 5:
                    correct_ans = opt
                    break
        
        if not correct_ans:
            for opt_raw in options_raw:
                if any(opt_raw.startswith(sym) for sym in ['g', 'iv', 'v', '©']):
                     correct_ans = re.sub(r'^[xgiv©*O)]+\s*', '', opt_raw).strip()
                     correct_ans = re.sub(r'\s+', ' ', correct_ans)
                     break

        questions.append({
            "text": q_text,
            "options": options,
            "correct": correct_ans,
            "explanation": explanation
        })

    # Output MD
    md_path = os.path.join(base_dir, "Life_Insurance_AZ_2026_MASTER_GUIDE.md")
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write("# 🏆 Life Insurance AZ 2026: Master Study Guide (Verified)\n\n")
        for i, q in enumerate(questions):
            f.write(f"## ❓ Question {i+1}\n**{q['text']}**\n\n### 📝 Options\n")
            if not q['options']: f.write("_Options missing_\n")
            for opt in q['options']:
                is_correct = (opt == q['correct'])
                f.write(f"- {'✅ [x]' if is_correct else '    [ ]'} {opt}\n")
            if q['explanation']:
                f.write(f"\n### 💡 Explanation\n> {q['explanation']}\n")
            f.write("\n---\n\n")

    # Output Word
    doc = Document()
    doc.add_heading('Life Insurance Master Study Guide', 0)
    for i, q in enumerate(questions):
        doc.add_heading(f'Question {i+1}', level=1)
        doc.add_paragraph(q['text']).bold = True
        for opt in q['options']:
            p = doc.add_paragraph(style='List Bullet')
            if opt == q['correct']:
                p.add_run(f"✅ [CORRECT] {opt}").bold = True
            else:
                p.add_run(opt)
        if q['explanation']:
            doc.add_paragraph(f"Explanation: {q['explanation']}").italic = True
    doc.save(os.path.join(base_dir, "Life_Insurance_Study_Guide.docx"))

    # Output PPT
    prs = Presentation()
    for i, q in enumerate(questions):
        # Slide 1: Question
        s1 = prs.slides.add_slide(prs.slide_layouts[1])
        s1.shapes.title.text = f"Question {i+1}"
        s1.placeholders[1].text = q['text']
        # Slide 2: Answer
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = f"Answer {i+1}"
        tf = s2.placeholders[1].text_frame
        for opt in q['options']:
            p = tf.add_paragraph()
            if opt == q['correct']:
                p.text = f"✅ {opt}"; p.font.bold = True
            else:
                p.text = f"☐ {opt}"
        if q['explanation']:
            p = tf.add_paragraph(); p.text = f"\nExplanation: {q['explanation'][:200]}..."
            p.font.size = PptPt(14); p.font.italic = True
    prs.save(os.path.join(base_dir, "Life_Insurance_Study_Guide.pptx"))

if __name__ == "__main__":
    rebuild()
